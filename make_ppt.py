from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

OUT_DIR = Path("outputs")
OUT_FILE = OUT_DIR / "dlo_2d_path_planning_method.pptx"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

TITLE_FONT = Pt(30)
BODY_FONT = Pt(18)
SMALL_FONT = Pt(15)
FORMULA_FONT = Pt(17)


def set_run(run, size=BODY_FONT, bold=False, font="Aptos"):
    run.font.name = font
    run.font.size = size
    run.font.bold = bold


def add_title(slide, title):
    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.32), Inches(12.2), Inches(0.55))
    frame = box.text_frame
    frame.clear()
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.LEFT
    run = paragraph.add_run()
    run.text = title
    set_run(run, TITLE_FONT, True)


def add_bullets(slide, x, y, w, h, items, font_size=BODY_FONT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = Inches(0.06)
    frame.margin_right = Inches(0.06)
    frame.margin_top = Inches(0.03)
    frame.margin_bottom = Inches(0.03)
    frame.clear()

    for idx, item in enumerate(items):
        paragraph = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        paragraph.text = item
        paragraph.level = 0
        paragraph.space_after = Pt(4)
        paragraph.font.size = font_size
        paragraph.font.name = "Aptos"


def add_lines(slide, x, y, w, h, lines, font_size=BODY_FONT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = Inches(0.08)
    frame.margin_right = Inches(0.08)
    frame.margin_top = Inches(0.04)
    frame.margin_bottom = Inches(0.04)
    frame.clear()

    for idx, line in enumerate(lines):
        paragraph = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        paragraph.space_after = Pt(4)
        run = paragraph.add_run()
        run.text = line
        set_run(run, font_size, False)


def add_section(slide, x, y, title, body_lines, width=5.9, height=1.9):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(height))
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = Inches(0.08)
    frame.margin_right = Inches(0.08)
    frame.margin_top = Inches(0.04)
    frame.margin_bottom = Inches(0.04)
    frame.clear()

    paragraph = frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = title
    set_run(run, BODY_FONT, True)
    paragraph.space_after = Pt(5)

    for line in body_lines:
        paragraph = frame.add_paragraph()
        paragraph.space_after = Pt(3)
        run = paragraph.add_run()
        run.text = line
        set_run(run, FORMULA_FONT, False)


def add_table(slide, x, y, w, h, rows, cols, data, font_size=SMALL_FONT):
    table_shape = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), Inches(h))
    table = table_shape.table
    for row_idx in range(rows):
        for col_idx in range(cols):
            cell = table.cell(row_idx, col_idx)
            cell.text = data[row_idx][col_idx]
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.alignment = PP_ALIGN.CENTER if col_idx == 0 else PP_ALIGN.LEFT
            for run in paragraph.runs:
                set_run(run, font_size, row_idx == 0)
    return table


def add_slide_1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Path Representation and Intermediate Shapes")

    add_section(
        slide,
        0.65,
        1.12,
        "Shape representation",
        [
            "DLO is discretized into N = 21 nodes.",
            "Each 2D shape is represented as X_j in R^(21 x 2).",
            "Given initial shape X_0 and target shape X_g.",
        ],
        width=5.95,
        height=1.65,
    )

    add_section(
        slide,
        6.85,
        1.12,
        "Planned path",
        [
            "X_0 -> X_1 -> ... -> X_M -> X_g",
            "Full planned path: X_path in R^(K x 21 x 2)",
        ],
        width=5.85,
        height=1.35,
    )

    add_section(
        slide,
        0.65,
        3.05,
        "Linear reference initialization",
        [
            "R_j = (1 - alpha_j) X_0 + alpha_j X_g",
            "alpha_j = j / (K - 1)",
        ],
        width=5.95,
        height=1.35,
    )

    add_section(
        slide,
        6.85,
        3.05,
        "Optimization role",
        [
            "The interpolation is only a reference.",
            "Final path is obtained by optimizing intermediate shapes",
            "under length and bending constraints.",
        ],
        width=5.85,
        height=1.65,
    )

    add_table(
        slide,
        1.6,
        5.38,
        10.1,
        0.8,
        2,
        3,
        [
            ["Variable", "Meaning", "Dimension"],
            ["X_j", "2D DLO shape at path index j", "21 x 2"],
        ],
    )


def add_slide_2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Intermediate-shape Optimization")

    add_section(
        slide,
        0.65,
        1.0,
        "Main objective",
        [
            "J_shape = w_ref J_ref + w_len J_len + w_bend J_bend + w_time J_time",
            "J_ref = sum_j ||X_j - R_j||^2",
            "J_len = sum_j sum_i (||p_(i+1)^j - p_i^j|| - Delta s)^2",
            "J_bend = sum_j ||D_2 X_j||^2",
            "J_time = sum_j ||X_(j+1) - X_j||^2",
        ],
        width=12.0,
        height=2.15,
    )

    add_table(
        slide,
        0.82,
        3.45,
        5.45,
        1.2,
        5,
        2,
        [
            ["Quantity", "Dimension / Definition"],
            ["X_j", "R^(21 x 2)"],
            ["D_1", "R^(20 x 21)"],
            ["D_2", "R^(19 x 21)"],
            ["Delta s", "L / (N - 1)"],
        ],
    )

    add_section(
        slide,
        6.65,
        3.28,
        "Implementation logic",
        [
            "1. Initialize X_path by linear interpolation.",
            "2. Smooth intermediate shapes.",
            "3. Penalize bending using D_2.",
            "4. Resample each shape by arc length.",
            "5. Keep X_0 and X_g fixed.",
            "6. Output planned_shapes.shape = (K, 21, 2).",
        ],
        width=5.95,
        height=2.75,
    )


def add_slide_3(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Three-force Inverse Actuation")

    add_section(
        slide,
        0.65,
        0.95,
        "Transition target",
        [
            "For each transition: X_j -> X_(j+1)",
            "delta q_j = vec(X_(j+1) - X_j) in R^42",
            "u_j = [s_1, s_2, s_3, F_x1, F_y1, F_x2, F_y2, F_x3, F_y3]",
            "Force spacing: |s_a - s_b| >= 0.01 m",
        ],
        width=12.0,
        height=1.8,
    )

    add_section(
        slide,
        0.65,
        2.85,
        "Force mapping and stiffness",
        [
            "F_node = A(s) f",
            "K delta q = A(s) f",
            "K = EA(D_1^T D_1) kron I_2 + EI(D_2^T D_2) kron I_2 + epsilon I",
            "Given force positions s: B(s) = K^-1 A(s)",
            "f* = (B^T B + lambda I)^-1 B^T delta q_j",
        ],
        width=7.35,
        height=2.35,
    )

    add_table(
        slide,
        8.28,
        2.9,
        4.35,
        1.45,
        5,
        2,
        [
            ["Quantity", "Dimension"],
            ["delta q_j", "R^42"],
            ["K", "R^(42 x 42)"],
            ["A(s)", "R^(42 x 6)"],
            ["f", "R^6"],
        ],
    )

    add_section(
        slide,
        8.28,
        4.65,
        "Final outputs",
        [
            "planned_shapes.npy",
            "force_actions.csv",
            "length_error.csv",
            "bending_energy.csv",
        ],
        width=4.35,
        height=1.5,
    )


def build_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    add_slide_1(prs)
    add_slide_2(prs)
    add_slide_3(prs)

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_FILE)
    return OUT_FILE


if __name__ == "__main__":
    output = build_presentation()
    print(f"Saved presentation to {output}")
